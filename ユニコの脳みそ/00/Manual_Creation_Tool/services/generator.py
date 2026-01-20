from pptx import Presentation
from pptx.util import Inches, Pt
import json
import os
import requests
from bs4 import BeautifulSoup
from openai import OpenAI
from jinja2 import Environment, FileSystemLoader

openai = None

def set_api_key(key):
    """
    API Keyを設定し、OpenAIクライアントを初期化する
    """
    global openai
    try:
        openai = OpenAI(api_key=key)
    except Exception as e:
        print(f"API Key Error: {e}")

def generate_dalle_image(prompt):
    """
    DALL-E 3を使って画像を生成し、URLを返す
    """
    if not openai: return None
    try:
        response = openai.images.generate(
            model="dall-e-3",
            prompt=prompt,
            size="1024x1024",
            quality="standard",
            n=1,
        )
        return response.data[0].url
    except Exception as e:
        print(f"DALL-E Error: {e}")
        return None

from urllib.parse import urljoin, urlparse
from duckduckgo_search import DDGS

def perform_company_research(company_name):
    """
    会社名から公式HPを検索し、Deep Scrapingを実行する
    """
    try:
        # 1. Search for Company URL
        with DDGS() as ddgs:
            results = list(ddgs.text(f"{company_name} 公式", max_results=1))
        
        if not results:
            return "Error: Company not found."
            
        target_url = results[0]['href']
        
        # 2. Deep Scrape the found URL
        scraped_content = scrape_website_content(target_url)
        
        return f"【Research Target: {company_name}】\n【Found URL: {target_url}】\n\n{scraped_content}"
        
    except Exception as e:
        return f"Research Error: {e}"
def scrape_website_content(url):
    """
    指定されたURLと、そのサイト内の主要ページ（最大3ページ）からテキストを抽出・統合する。
    """
    try:
        headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'}
        
        # 1. Main Page Fetch
        session = requests.Session()
        resp = session.get(url, headers=headers, timeout=10)
        resp.encoding = resp.apparent_encoding
        soup = BeautifulSoup(resp.text, 'html.parser')
        
        # Extract Main Content
        full_text = f"【Origin URL: {url}】\n"
        full_text += _extract_clean_text(soup) + "\n\n"
        
        # 2. Find Internal Links
        domain = urlparse(url).netloc
        links = []
        for a in soup.find_all('a', href=True):
            href = a['href']
            full_url = urljoin(url, href)
            parsed = urlparse(full_url)
            
            # Internal link check (same domain, not anchor, not file)
            if parsed.netloc == domain and "#" not in href and not href.lower().endswith(('.pdf', '.jpg', '.png')):
                # Priorities: Company, Service, About, Recruit
                score = 1 # Base score for all internal links
                if any(k in href.lower() for k in ['about', 'company', 'profile', 'gaiyo']): score += 5
                if any(k in href.lower() for k in ['service', 'business', 'works', 'product']): score += 3
                if any(k in href.lower() for k in ['recruit', 'job', 'career']): score += 2
                
                # Include ALL internal links (removed score > 0 check)
                if full_url not in [l[1] for l in links] and full_url != url:
                    links.append((score, full_url))
        
        # Sort by score and take top 15 (Capture Everything Strategy)
        links.sort(key=lambda x: x[0], reverse=True)
        target_links = [l[1] for l in links[:15]]
        
        # 3. Scrape Sub Pages
        for sub_url in target_links:
            try:
                sub_resp = session.get(sub_url, headers=headers, timeout=5)
                sub_resp.encoding = sub_resp.apparent_encoding
                sub_soup = BeautifulSoup(sub_resp.text, 'html.parser')
                
                full_text += f"【Sub Page: {sub_url}】\n"
                full_text += _extract_clean_text(sub_soup) + "\n\n"
            except:
                continue
                
        return full_text[:50000] # Increased limit to 50k chars for massive context
        
    except Exception as e:
        return f"Scraping Error: {e}"

def extract_text_from_pdf(file_obj):
    """
    PDFファイルからテキストを抽出する
    """
    try:
        from pypdf import PdfReader
        reader = PdfReader(file_obj)
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
        return text
    except Exception as e:
        return f"PDF Error: {e}"

def extract_text_from_pptx(file_obj):
    """
    PPTXファイルからテキストを抽出する
    """
    try:
        from pptx import Presentation
        prs = Presentation(file_obj)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        return f"PPTX Error: {e}"

def _extract_clean_text(soup):
    for script in soup(["script", "style", "nav", "footer", "header", "iframe", "noscript"]):
        script.decompose()
    text = soup.get_text()
    lines = (line.strip() for line in text.splitlines())
    chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
    return '\n'.join(chunk for chunk in chunks if chunk)

def generate_hearing_questions(input_text):
    """
    【Hearing Bot】
    入力されたテキスト（メモ）の「不明点」「暗黙知」を洗い出し、
    ユーザーに確認すべき3つの質問を生成する。
    """
    if not openai: return ["Error: OpenAI not installed"]
    
    prompt = f"""
    あなたは優秀なビジネスアナリストです。
    以下の「マニュアル化したい業務メモ」を読み、
    マニュアルを作成する上で「情報が足りない点」「曖昧な点（暗黙知）」を3つ特定してください。
    
    【入力メモ】
    {input_text}
    
    出力は以下のJSON形式のみにしてください。
    {{"q1": "質問1", "q2": "質問2", "q3": "質問3"}}
    """
    
    try:
        response = openai.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "user", "content": prompt}],
            response_format={"type": "json_object"},
            temperature=0.7
        )
        import json
        return json.loads(response.choices[0].message.content)
    except Exception as e:
        return {"q1": f"Error: {e}", "q2": "", "q3": ""}

def generate_structure_options_with_hearing(input_text, hearing_answers, mode):
    """
    Hearingの回答を含めて構成案を生成する
    """
    # 質問と回答を統合
    enhanced_input = f"""
    【元のメモ】
    {input_text}
    
    【追加ヒアリング情報】
    Q1の回答: {hearing_answers.get('a1', '')}
    Q2の回答: {hearing_answers.get('a2', '')}
    Q3の回答: {hearing_answers.get('a3', '')}
    """
    return generate_structure_options(enhanced_input, mode)

# ... (generate_structure_options - minor mod to accept prompt directly or reusing logic)

def generate_structure_options(input_text, mode):
    # (Same as before but keep it robust)
    if not openai: return ["Error: OpenAI not installed"]
    
    prompts = {
        "Manual": "業務マニュアルの目次構成案",
        # ... (rest same)
    }
    target = prompts.get(mode, "構成案")
    
    prompt = f"""
    あなたはプロのドキュメント作成コンサルタントです。
    以下の情報を元に、{target} を **3パターン** 提案してください。
    
    【入力情報】
    {input_text}
    
    【出力条件】
    - Option 1: 基本的・網羅的な構成
    - Option 2: 初心者向けの分かりやすさ重視の構成
    - Option 3: 要点重視・スピード重視の簡潔な構成
    - 各案は「タイトル」と「概要」を明確に分けて記述してください。
    """
    
    try:
        response = openai.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.7
        )
        return response.choices[0].message.content
    except Exception as e:
        return f"Error: {e}"

def generate_slide_content_json(final_text):
    """
    最終テキストを「スライド用JSONデータ」に変換する
    """
    prompt = f"""
    以下のマニュアルテキストを、PowerPointスライド（5〜10枚）にするためのJSONデータに変換してください。
    
    【テキスト】
    {final_text}
    
    【JSON形式】
    {{
        "slides": [
            {{"title": "スライド1タイトル", "points": ["箇条書き1", "箇条書き2"]}},
            ...
        ]
    }}
    """
    try:
        response = openai.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "user", "content": prompt}],
            response_format={"type": "json_object"},
            temperature=0.7
        )
        import json
        return json.loads(response.choices[0].message.content)
    except:
        return {"slides": []}

def create_presentation(slide_data):
    """
    JSONデータから実際のpptxファイルを生成する
    """
    prs = Presentation()
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "業務マニュアル"
    slide.placeholders[1].text = "Generated by AI Manual Architect"
    
    # Content Slides
    bullet_slide_layout = prs.slide_layouts[1]
    
    for s_data in slide_data.get("slides", []):
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        shapes.title.text = s_data.get("title", "No Title")
        
        tf = shapes.placeholders[1].text_frame
        for point in s_data.get("points", []):
            p = tf.add_paragraph()
            p.text = point
            p.font.size = Pt(18)
            
    filename = "generated_manual.pptx"
    prs.save(filename)
    return filename


# ... (Previous imports)

# --- CAST DEFINITIONS (Bible Compliant) ---
CAST_PROMPTS = {
    "Standard": {
        "Manabu": "Young male rookie (22), short messy black hair, ill-fitting navy suit, expressive face.",
        "Shigoto": "Professional female mentor (28), glasses, ponytail, sharp beige suit, calm expression.",
        "Style": "Japanese commercial manga style, black and white line art, clean lines."
    },
    "Factory": {
        "Manabu": "Young factory worker (20), wearing safety helmet and green work uniform, safety gloves.",
        "Shigoto": "Veteran craftsman (55), stubble, towel around neck, grey work uniform, stern but kind eyes.",
        "Style": "Seinen manga style, detailed machinery background, gritty but clean."
    },
    "Sales": {
        "Manabu": "Eager sales rookie (24), perfectly gelled hair, sharp slim suit, bright smile but nervous eyes.",
        "Shigoto": "Elite Sales Manager (35), slicked back hair, expensive 3-piece suit, aura of confidence.",
        "Style": "Business manga style, corporate background, sharp angles."
    }
}


# ... (Previous imports)

def generate_html_content(selected_option, final_text, mode, cast_key="Standard"):
    """
    HTMLを生成するメイン関数
    """
    base_dir = os.path.dirname(os.path.abspath(__file__))
    template_dir = os.path.join(base_dir, 'templates')
    env = Environment(loader=FileSystemLoader(template_dir))
    
    if mode == "Manga":
        template = env.get_template('manga.html')
        data = parse_manga_content(final_text, cast_key)
    elif mode == "Visual":
        template = env.get_template('visual.html')
        data = parse_visual_content(final_text)
    elif mode == "LP":
        template = env.get_template('lp.html')
        data = parse_lp_content(final_text)
    else: # Business
        template = env.get_template('business.html')
        try:
            import markdown
            html_body = markdown.markdown(final_text, extensions=['extra'])
        except:
            html_body = final_text.replace("\n", "<br>")
        data = {"title": "業務マニュアル", "content": html_body}

    return render_html_from_data(data, mode)

def render_html_from_data(data, mode):
    """
    データ辞書からHTMLをレンダリングする (画像挿入後の再描画用)
    """
    base_dir = os.path.dirname(os.path.abspath(__file__))
    template_dir = os.path.join(base_dir, 'templates')
    env = Environment(loader=FileSystemLoader(template_dir))
    
    if mode == "Manga":
        template = env.get_template('manga.html')
    elif mode == "Visual":
        template = env.get_template('visual.html')
    elif mode == "LP":
        template = env.get_template('lp.html')
    else:
        template = env.get_template('business.html')
        
    return template.render(data)

def generate_thumbnail_html(final_text):
    """
    サムネイル用HTMLを生成する
    """
    base_dir = os.path.dirname(os.path.abspath(__file__))
    template_dir = os.path.join(base_dir, 'templates')
    env = Environment(loader=FileSystemLoader(template_dir))
    data = parse_thumbnail_content(final_text)
    return render_thumbnail_from_data(data)

def render_thumbnail_from_data(data):
    """
    サムネイルデータからHTMLをレンダリングする
    """
    base_dir = os.path.dirname(os.path.abspath(__file__))
    template_dir = os.path.join(base_dir, 'templates')
    env = Environment(loader=FileSystemLoader(template_dir))
    template = env.get_template('thumbnail.html')
    return template.render(data)

def parse_thumbnail_content(text):
    """
    マニュアル内容から「サムネイル用のキャッチコピー」と「デザイン」を抽出・決定する
    """
    if not openai: 
        return {
            "main_text": "AI Manual", 
            "sub_text": "Generator", 
            "bg_style": "linear-gradient(135deg, #1e3c72 0%, #2a5298 100%)"
        }
    
    prompt = f"""
    以下のマニュアル内容を元に、YouTubeサムネイル風の「極太キャッチコピー」を作成してください。
    内容の「緊急度」や「雰囲気」に合わせて、背景色も指定してください。
    
    【内容】
    {text}
    
    【JSONフォーマット】
    {{
        "main_text": "一番目立つ文字（例：閲覧注意！, 〇〇の真実, 3分でわかる）",
        "sub_text": "補足テキスト（例：新人必見, 保存版, 業務改革）",
        "bg_type": "Design Theme (Warning / Trust / Premium / Pop)",
        "img_prompt": "サムネイルの背景として相応しい画像の生成プロンプト (例: futuristic office, dark hacker room)"
    }}
    """
    try:
        res = openai.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "user", "content": prompt}],
            response_format={"type": "json_object"},
            temperature=0.8
        )
        data = json.loads(res.choices[0].message.content)
        
        # Theme to CSS Mapping
        themes = {
            "Warning": "linear-gradient(45deg, #ff0000, #330000)", # Panic Red
            "Trust": "linear-gradient(135deg, #00c6ff, #0072ff)", # Trust Blue
            "Premium": "linear-gradient(to right, #bf953f, #fcf6ba, #b38728, #fbf5b7, #aa771c)", # Gold
            "Pop": "linear-gradient(120deg, #f6d365 0%, #fda085 100%)", # Orange/Yellow
            "Default": "linear-gradient(135deg, #667eea 0%, #764ba2 100%)" # Purple
        }
        
        bg_style = themes.get(data.get("bg_type", "Default"), themes["Default"])
        
        return {
            "main_text": data.get("main_text", "Manual"),
            "sub_text": data.get("sub_text", "Update"),
            "bg_style": bg_style,
            "img_prompt": data.get("img_prompt", "")
        }
    except:
        return {"main_text": "Error", "sub_text": "Try Again", "bg_style": "black"}

# ... (parse_manga_content and others remain)

def parse_manga_content(text, cast_key="Standard"):
    """
    Manga Mode Parsing with Cast Injection
    """
    if not openai: return {"title": "Error", "sections": []}
    
    cast = CAST_PROMPTS.get(cast_key, CAST_PROMPTS["Standard"])
    
    prompt = f"""
    以下のマニュアル内容を、マンガ形式のJSONデータに変換してください。
    
    【内容】
    {text}
    
    【キャラクター設定 (Cast: {cast_key})】
    - Rookie: {cast['Manabu']}
    - Mentor: {cast['Shigoto']}
    - Art Style: {cast['Style']}
    
    【JSONフォーマット】
    {{
        "title": "タイトル",
        "sections": [
            {{
                "narrative": "...",
                "img_prompt": "Generate a manga panel... [Include Character Descriptions Here] ...", 
                "dialogues": [ {{"character": "Rookie", "text": "..."}}, {{"character": "Mentor", "text": "..."}} ]
            }}
        ]
    }}
    """
    try:
        res = openai.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "user", "content": prompt}],
            response_format={"type": "json_object"},
            temperature=0.7
        )
        return json.loads(res.choices[0].message.content)
    except:
        return {"title": "Parse Error", "sections": []}

def parse_visual_content(text):
    """
    Visual Mode with Smart Threshold Logic
    """
    if not openai: return {"title": "Error", "sections": []}
    
    prompt = f"""
    マニュアル内容をVisualテンプレート用JSONに変換してください。
    【重要】chart_dataを作成する際、'backgroundColor'をデータの値に基づいて動的に設定してください。
    例: 値が悪い(高い/低い)場合は 'rgba(255, 99, 132, 0.6)' (赤)、良い場合は 'rgba(75, 192, 192, 0.6)' (緑) など。
    
    【内容】
    {text}
    
    【JSON出力】
    (Standard Visual JSON format, ensure backgroundColor is an array matching data length)
    """
    try:
        res = openai.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "user", "content": prompt}],
            response_format={"type": "json_object"},
            temperature=0.7
        )
        return json.loads(res.choices[0].message.content)
    except:
        return {"title": "Parse Error", "sections": []}

def parse_lp_content(text):
    """
    LP Content Parsing (PAS Model)
    """
    if not openai: return {"hero": {"headline": "Error"}}
    
    prompt = f"""
    以下のサービス・業務内容を元に、高成約率のランディングページ(LP)構成を作成してください。
    PASモデル（Problem, Agitation, Solution）を使用すること。
    
    【入力内容】
    {text}
    
    【JSONフォーマット】
    {{
        "title": "ページタイトル",
        "hero": {{ "headline": "キャッチコピー", "subheadline": "サブコピー", "cta_text": "ボタン文言" }},
        "pas": {{ "problem": "問題提起", "agitation": "問題の深堀り（煽り）", "solution": "解決策" }},
        "features": [ {{ "title": "特徴1", "desc": "説明" }}, {{ "title": "特徴2", "desc": "説明" }}, {{ "title": "特徴3", "desc": "説明" }} ],
        "testimonial": {{ "text": "お客様の声", "author": "誰某" }},
        "cta": {{ "headline": "クロージングコピー", "subheadline": "最後の一押し", "button_text": "今すぐ申し込む" }}
    }}
    """
    try:
        res = openai.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "user", "content": prompt}],
            response_format={"type": "json_object"},
            temperature=0.7
        )
        return json.loads(res.choices[0].message.content)
    except:
        return {"hero": {"headline": "Error"}}

# ... (rest of file)


def generate_final_content(selected_option, input_text, mode):
    """
    選ばれた構成案を元に、最終成果物（Markdown詳細）を生成する。
    """
    if not openai: return "Error"
    
    system_prompts = {
        "Manual": "あなたは「誰でも読める分かりやすい業務マニュアル」のライターです。Markdownで見出し、手順、注意点を構造化して記述してください。",
        "Flowchart": "あなたは「Mermaid記法」の専門家です。業務フローをMermaid形式のコードブロック（graph TD）で出力してください。テキスト解説も添えてください。",
        "FAQ": "あなたは「カスタマーサポートのプロ」です。Q&A形式で、質問と、共感を含めた丁寧な回答を作成してください。",
        "Slide": "あなたは「プロのプレゼン資料作成者」です。各スライドの「タイトル」「本文（箇条書き）」「スピーカーノート」を構成してください。"
    }
    
    prompt = f"""
    以下の「採用された構成案」と「元の入力情報」を元に、
    最終的な成果物を作成してください。
    
    【採用された構成案】
    {selected_option}
    
    【元の入力情報】
    {input_text}
    
    出力は高品質な日本語で、そのままコピペして使えるレベルに仕上げてください。
    """
    
    try:
        response = openai.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": system_prompts.get(mode, "You are a professional writer.")},
                {"role": "user", "content": prompt}
            ],
            temperature=0.7
        )
        return response.choices[0].message.content
    except Exception as e:
        return f"Error: {e}"

def generate_quote_pdf():
    """
    松・竹・梅の3プラン見積もりPDFを作成する
    """
    try:
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import A4
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont
        import os
        
        filename = "Quote_v1.pdf"
        p = canvas.Canvas(filename, pagesize=A4)
        width, height = A4
        
        # Font Registration (Try common Windows fonts)
        font_name = "Helvetica"
        try:
            # Check for generic gothic font usually present in Windows
            font_path = "C:\\Windows\\Fonts\\msgothic.ttc" 
            if os.path.exists(font_path):
                pdfmetrics.registerFont(TTFont('Gothic', font_path))
                font_name = 'Gothic'
            else:
                 # Fallback or try another
                 pass
        except:
            pass
            
        # Draw Header
        p.setFont(font_name, 24)
        p.drawString(50, height - 80, "お見積書 (Quotation)")
        
        p.setFont(font_name, 12)
        p.drawString(50, height - 120, "件名: AIマニュアル作成代行サービスのご提案")
        p.drawString(50, height - 140, "発行日: 2026年1月15日")
        p.drawString(400, height - 120, "株式会社Unico Brain") # Generic Sender
        
        # Draw Plans
        y = height - 200
        
        plans = [
             {"name": "【松】Premium Plan (Agency)", "price": "¥300,000", "desc": "マンガLP + 動的グラフ + 30Pマニュアル + 修正無制限"},
             {"name": "【竹】Standard Plan (Biz)", "price": "¥100,000", "desc": "HTMLマニュアル + 標準デザイン + 10P + 修正2回"},
             {"name": "【梅】Light Plan (Text)", "price": "¥50,000", "desc": "テキスト調整のみ + PDF納品"}
        ]
        
        for plan in plans:
            p.setStrokeColorRGB(0.2, 0.2, 0.2)
            p.rect(50, y - 80, 500, 80, fill=0)
            
            p.setFont(font_name, 16)
            p.drawString(70, y - 30, plan["name"])
            
            p.setFont(font_name, 18)
            p.drawString(400, y - 30, plan["price"])
            
            p.setFont(font_name, 10)
            p.drawString(70, y - 60, plan["desc"])
            
            y -= 100
            
        # Footer
        p.setFont(font_name, 10)
        p.drawString(50, 100, "※本見積書の有効期限は発行から2週間とします。")
        
        p.save()
        return filename
    except Exception as e:
        print(f"PDF Gen Error: {e}")
        return None
